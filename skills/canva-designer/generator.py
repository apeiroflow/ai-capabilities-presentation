"""
Canva Slide Generator
Creates Canva-branded HTML presentations using Reveal.js.
Follows Canva Connect brand guidelines: https://www.canva.dev/docs/connect/guidelines/brand/
"""
import html
import os
import json
from typing import List, Dict, Optional, Union


class CanvaSlideGenerator:
    """Generate Canva-branded HTML presentations."""

    BRAND_COLORS = {
        "purple": "#8B3DFF",
        "purple_dark": "#6D28D9",
        "teal": "#00C4CC",
        "teal_dark": "#0099A8",
        "dark": "#1A1A2E",
        "navy": "#16213E",
        "grey": "#6B7280",
        "light_grey": "#E5E7EB",
        "off_white": "#F9FAFB",
        "white": "#FFFFFF",
        "green": "#10B981",
        "amber": "#F59E0B",
        "pink": "#EC4899",
    }

    BRAND_COLORS_RGB = {
        "purple": (139, 61, 255),
        "teal": (0, 196, 204),
        "dark": (26, 26, 46),
        "grey": (107, 114, 128),
        "white": (255, 255, 255),
        "green": (16, 185, 129),
    }

    FONT_URL = "https://fonts.googleapis.com/css2?family=Syne:wght@400;600;700;800&family=DM+Sans:wght@300;400;500;600&display=swap"
    FONT_DISPLAY = "'Syne', sans-serif"
    FONT_BODY = "'DM Sans', sans-serif"

    def __init__(self, title: str = "Canva Presentation"):
        self.title = title
        self._slides_html: List[str] = []
        self._slides_data: List[Dict] = []

    # ─── Security helpers ────────────────────────────────────────────────────

    @staticmethod
    def _e(text) -> str:
        """Escape a single string for safe HTML insertion."""
        return html.escape(str(text)) if text is not None else ""

    @staticmethod
    def _el(items: List) -> List[str]:
        """Escape a list of strings."""
        return [html.escape(str(i)) for i in items]

    @staticmethod
    def _ed(d: Dict) -> Dict:
        """Escape all string values in a dict."""
        return {k: html.escape(str(v)) if isinstance(v, str) else v for k, v in d.items()}

    # ─── Slide builders ──────────────────────────────────────────────────────

    def add_title_slide(
        self,
        title: str,
        subtitle: str = "",
        author: str = "",
        badge: str = "",
        show_powered_by: bool = True,
    ) -> "CanvaSlideGenerator":
        t, s, a, b = self._e(title), self._e(subtitle), self._e(author), self._e(badge)
        badge_html = f'<span class="badge">{b}</span>' if b else ""
        subtitle_html = f'<p class="slide-subtitle">{s}</p>' if s else ""
        author_html = f'<p class="author">{a}</p>' if a else ""
        powered = '<div class="powered-by"><span class="powered-text">Powered by</span><span class="canva-wordmark">Canva</span></div>' if show_powered_by else ""
        slide = f"""
<section class="slide-title">
  <div class="title-bg-mesh"></div>
  <div class="center-container">
    {badge_html}
    <h1 class="title-heading">{t}</h1>
    <div class="accent-bar"></div>
    {subtitle_html}
    {author_html}
    {powered}
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "title", "title": title, "subtitle": subtitle})
        return self

    def add_section_slide(
        self,
        title: str,
        subtitle: str = "",
        color: str = "purple",
    ) -> "CanvaSlideGenerator":
        t, s = self._e(title), self._e(subtitle)
        sub_html = f'<p class="section-sub">{s}</p>' if s else ""
        bg = self.BRAND_COLORS.get(color, self.BRAND_COLORS["purple"])
        slide = f"""
<section class="slide-section" style="background: linear-gradient(135deg, {bg} 0%, {self.BRAND_COLORS['dark']} 100%);">
  <div class="section-noise"></div>
  <div class="center-container">
    <h2 class="section-title">{t}</h2>
    {sub_html}
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "section", "title": title})
        return self

    def add_content_slide(
        self,
        title: str,
        content: List[str],
        badge: str = "",
        subtitle: str = "",
    ) -> "CanvaSlideGenerator":
        t, b, s = self._e(title), self._e(badge), self._e(subtitle)
        items = "".join(f"<li>{self._e(i)}</li>" for i in content)
        badge_html = f'<span class="badge badge-outline">{b}</span>' if b else ""
        sub_html = f'<p class="slide-subtitle small">{s}</p>' if s else ""
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    {sub_html}
    <ul class="content-list">{items}</ul>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "content", "title": title, "items": content})
        return self

    def add_stat_slide(
        self,
        title: str,
        stats: List[Dict],
        badge: str = "",
        subtitle: str = "",
        columns: int = 3,
    ) -> "CanvaSlideGenerator":
        """stats: list of {number, label, icon?, sublabel?}"""
        t, b, s = self._e(title), self._e(badge), self._e(subtitle)
        badge_html = f'<span class="badge">{b}</span>' if b else ""
        sub_html = f'<p class="slide-subtitle small">{s}</p>' if s else ""
        col_cls = {2: "two", 3: "", 4: "four"}.get(columns, "")
        boxes = ""
        for st in stats:
            d = self._ed(st)
            icon = f'<div class="stat-icon">{d.get("icon","")}</div>' if d.get("icon") else ""
            sub = f'<div class="stat-sublabel">{d.get("sublabel","")}</div>' if d.get("sublabel") else ""
            boxes += f"""<div class="stat-box"><{icon}<div class="stat-value">{d.get("number","")}</div><div class="stat-label">{d.get("label","")}</div>{sub}</div>"""
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    {sub_html}
    <div class="stat-grid {col_cls}">{boxes}</div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "stat", "title": title, "stats": stats})
        return self

    def add_hero_stat_slide(
        self,
        title: str,
        number: str,
        label: str,
        sublabel: str = "",
        badge: str = "",
    ) -> "CanvaSlideGenerator":
        t, n, l, sl, b = self._e(title), self._e(number), self._e(label), self._e(sublabel), self._e(badge)
        badge_html = f'<span class="badge">{b}</span>' if b else ""
        sub_html = f'<p class="hero-sublabel">{sl}</p>' if sl else ""
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    <div class="hero-stat">
      <div class="hero-number">{n}</div>
      <div class="hero-label">{l}</div>
      {sub_html}
    </div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "hero_stat", "title": title, "number": number})
        return self

    def add_flow_slide(
        self,
        title: str,
        steps: List[str],
        badge: str = "",
        subtitle: str = "",
        accent_last: bool = True,
    ) -> "CanvaSlideGenerator":
        t, b, s = self._e(title), self._e(badge), self._e(subtitle)
        badge_html = f'<span class="badge">{b}</span>' if b else ""
        sub_html = f'<p class="slide-subtitle small">{s}</p>' if s else ""
        steps_html = ""
        for i, step in enumerate(steps):
            is_last = i == len(steps) - 1
            accent = "accent" if (accent_last and is_last) else ""
            arrow = '<div class="flow-arrow">→</div>' if not is_last else ""
            steps_html += f'<div class="flow-step {accent}"><div class="flow-num">{i+1}</div><div class="flow-text">{self._e(step)}</div></div>{arrow}'
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    {sub_html}
    <div class="flow-container">{steps_html}</div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "flow", "title": title, "steps": steps})
        return self

    def add_bar_chart_slide(
        self,
        title: str,
        data: List[Dict],
        badge: str = "",
        subtitle: str = "",
    ) -> "CanvaSlideGenerator":
        """data: list of {label, value, max?}"""
        t, b, s = self._e(title), self._e(badge), self._e(subtitle)
        badge_html = f'<span class="badge">{b}</span>' if b else ""
        sub_html = f'<p class="slide-subtitle small">{s}</p>' if s else ""
        max_val = max((d.get("value", 0) for d in data), default=100)
        bars = ""
        for item in data:
            d = self._ed(item)
            pct = round((item.get("value", 0) / max_val) * 100)
            bars += f"""<div class="bar-row"><div class="bar-label">{d.get("label","")}</div><div class="bar-track"><div class="bar-fill" style="width:{pct}%"></div></div><div class="bar-value">{d.get("value","")}</div></div>"""
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    {sub_html}
    <div class="bar-chart">{bars}</div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "bar_chart", "title": title, "data": data})
        return self

    def add_comparison_slide(
        self,
        title: str,
        left: Dict,
        right: Dict,
        badge: str = "",
    ) -> "CanvaSlideGenerator":
        """left/right: {title, items: [str]}"""
        t, b = self._e(title), self._e(badge)
        badge_html = f'<span class="badge badge-outline">{b}</span>' if b else ""

        def side(d, cls):
            items = "".join(f"<li>{self._e(i)}</li>" for i in d.get("items", []))
            return f'<div class="comparison-side {cls}"><h3>{self._e(d.get("title",""))}</h3><ul>{items}</ul></div>'

        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    <div class="comparison">{side(left, "left")}{side(right, "right")}</div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "comparison", "title": title})
        return self

    def add_two_col_slide(
        self,
        title: str,
        left_content: str,
        right_content: str,
        badge: str = "",
    ) -> "CanvaSlideGenerator":
        t, b = self._e(title), self._e(badge)
        badge_html = f'<span class="badge badge-outline">{b}</span>' if b else ""
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    <div class="two-col">
      <div class="col">{left_content}</div>
      <div class="col">{right_content}</div>
    </div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "two_col", "title": title})
        return self

    def add_feature_grid_slide(
        self,
        title: str,
        features: List[Dict],
        badge: str = "",
        columns: int = 2,
    ) -> "CanvaSlideGenerator":
        """features: list of {title, description, icon?}"""
        t, b = self._e(title), self._e(badge)
        badge_html = f'<span class="badge">{b}</span>' if b else ""
        cards = ""
        for f in features:
            d = self._ed(f)
            icon = f'<div class="feat-icon">{d.get("icon","")}</div>' if d.get("icon") else ""
            cards += f'<div class="feature-card">{icon}<h4>{d.get("title","")}</h4><p>{d.get("description","")}</p></div>'
        grid_style = f'grid-template-columns: repeat({columns}, 1fr);'
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    <div class="feature-grid" style="{grid_style}">{cards}</div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "feature_grid", "title": title})
        return self

    def add_quote_slide(
        self,
        quote: str,
        author: str = "",
        role: str = "",
        badge: str = "",
    ) -> "CanvaSlideGenerator":
        q, a, r, b = self._e(quote), self._e(author), self._e(role), self._e(badge)
        badge_html = f'<span class="badge">{b}</span>' if b else ""
        attr = f'<div class="quote-author">{a}{", " + r if r else ""}</div>' if a else ""
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <div class="quote-block">
      <div class="quote-mark">"</div>
      <blockquote>{q}</blockquote>
      {attr}
    </div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "quote", "quote": quote})
        return self

    def add_badge_slide(
        self,
        badge_label: str,
        title: str,
        bullets: List[str],
        subtitle: str = "",
    ) -> "CanvaSlideGenerator":
        bl, t, s = self._e(badge_label), self._e(title), self._e(subtitle)
        items = "".join(f"<li>{self._e(i)}</li>" for i in bullets)
        sub_html = f'<p class="slide-subtitle small">{s}</p>' if s else ""
        slide = f"""
<section>
  <div class="center-container">
    <span class="badge badge-large">{bl}</span>
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    {sub_html}
    <ul class="content-list">{items}</ul>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "badge", "title": title})
        return self

    def add_phase_cards_slide(
        self,
        title: str,
        phases: List[Dict],
        badge: str = "",
    ) -> "CanvaSlideGenerator":
        """phases: list of {title, description, label?}"""
        t, b = self._e(title), self._e(badge)
        badge_html = f'<span class="badge">{b}</span>' if b else ""
        cards = ""
        for i, ph in enumerate(phases):
            d = self._ed(ph)
            lbl = d.get("label", f"Phase {i+1}")
            cards += f'<div class="phase-card"><div class="phase-num">{i+1}</div><div class="phase-label">{lbl}</div><h4>{d.get("title","")}</h4><p>{d.get("description","")}</p></div>'
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    <div class="phase-grid">{cards}</div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "phase_cards", "title": title})
        return self

    def add_timeline_slide(
        self,
        title: str,
        milestones: List[Dict],
        badge: str = "",
    ) -> "CanvaSlideGenerator":
        """milestones: list of {date, title, description?}"""
        t, b = self._e(title), self._e(badge)
        badge_html = f'<span class="badge">{b}</span>' if b else ""
        items = ""
        for i, m in enumerate(milestones):
            d = self._ed(m)
            desc = f'<p>{d.get("description","")}</p>' if d.get("description") else ""
            side = "left" if i % 2 == 0 else "right"
            items += f'<div class="timeline-item {side}"><div class="timeline-dot"></div><div class="timeline-content"><div class="timeline-date">{d.get("date","")}</div><h4>{d.get("title","")}</h4>{desc}</div></div>'
        slide = f"""
<section>
  <div class="center-container">
    {badge_html}
    <h2>{t}</h2>
    <div class="accent-bar"></div>
    <div class="timeline">{items}</div>
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "timeline", "title": title})
        return self

    def add_thank_you_slide(
        self,
        title: str = "Thank You",
        subtitle: str = "",
        contact: Optional[Dict] = None,
        show_powered_by: bool = True,
    ) -> "CanvaSlideGenerator":
        t, s = self._e(title), self._e(subtitle)
        sub_html = f'<p class="slide-subtitle">{s}</p>' if s else ""
        contact_html = ""
        if contact:
            rows = "".join(f'<p><strong>{self._e(k)}:</strong> {self._e(v)}</p>' for k, v in contact.items())
            contact_html = f'<div class="contact-box">{rows}</div>'
        powered = '<div class="powered-by large"><span class="powered-text">Powered by</span><span class="canva-wordmark">Canva</span></div>' if show_powered_by else ""
        slide = f"""
<section class="slide-title">
  <div class="title-bg-mesh"></div>
  <div class="center-container">
    <h1 class="title-heading">{t}</h1>
    <div class="accent-bar"></div>
    {sub_html}
    {contact_html}
    {powered}
  </div>
</section>"""
        self._slides_html.append(slide)
        self._slides_data.append({"type": "thank_you", "title": title})
        return self

    # ─── HTML generation ─────────────────────────────────────────────────────

    def _build_css(self) -> str:
        c = self.BRAND_COLORS
        return f"""
        <style>
        @import url('{self.FONT_URL}');

        :root {{
            --purple: {c['purple']};
            --purple-dark: {c['purple_dark']};
            --teal: {c['teal']};
            --teal-dark: {c['teal_dark']};
            --dark: {c['dark']};
            --navy: {c['navy']};
            --grey: {c['grey']};
            --light-grey: {c['light_grey']};
            --off-white: {c['off_white']};
            --white: {c['white']};
            --green: {c['green']};
            --amber: {c['amber']};
            --pink: {c['pink']};
        }}

        * {{ box-sizing: border-box; }}

        body {{ background: var(--off-white); }}

        .reveal {{
            font-family: {self.FONT_BODY};
            background: var(--off-white);
        }}

        .reveal .slides section {{
            padding: 40px 60px;
            height: 100%;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
        }}

        .reveal h1, .reveal h2, .reveal h3, .reveal h4 {{
            font-family: {self.FONT_DISPLAY};
            text-align: center !important;
            text-transform: none;
            letter-spacing: -0.02em;
        }}

        .reveal h1 {{ color: var(--dark); font-size: 2.4em; font-weight: 800; line-height: 1.1; margin-bottom: 0.2em; }}
        .reveal h2 {{ color: var(--dark); font-size: 1.6em; font-weight: 700; margin-bottom: 0.2em; }}
        .reveal h3 {{ color: var(--dark); font-size: 1.0em; font-weight: 600; }}
        .reveal h4 {{ color: var(--dark); font-size: 0.85em; font-weight: 600; margin-bottom: 8px; }}
        .reveal p  {{ color: var(--dark); font-size: 0.85em; line-height: 1.6; text-align: center; }}

        .reveal ul {{ list-style: none; text-align: left; margin: 0; padding: 0; }}
        .reveal ul li {{
            margin-bottom: 10px; padding-left: 28px; font-size: 0.8em;
            line-height: 1.4; position: relative; color: var(--dark);
        }}
        .reveal ul li::before {{
            content: "";
            position: absolute; left: 0; top: 7px;
            width: 9px; height: 9px;
            background: linear-gradient(135deg, var(--purple), var(--teal));
            border-radius: 50%;
        }}

        /* ── Layout ── */
        .center-container {{
            display: flex; flex-direction: column; align-items: center;
            text-align: center; width: 100%; max-width: 1100px;
        }}

        .two-col {{ display: grid; grid-template-columns: 1fr 1fr; gap: 28px; width: 100%; max-width: 1100px; margin: 20px auto; }}
        .col {{ display: flex; flex-direction: column; align-items: center; }}

        /* ── Accent bar ── */
        .accent-bar {{
            height: 3px;
            background: linear-gradient(90deg, var(--purple) 0%, var(--teal) 100%);
            width: 72px; margin: 12px auto 20px; border-radius: 2px;
        }}

        /* ── Badge ── */
        .badge {{
            display: inline-block;
            background: linear-gradient(135deg, var(--purple) 0%, var(--teal) 100%);
            color: white; padding: 5px 18px; border-radius: 30px;
            font-size: 0.62em; font-weight: 600; text-transform: uppercase;
            letter-spacing: 1.2px; margin-bottom: 14px; font-family: {self.FONT_BODY};
        }}
        .badge-outline {{
            background: transparent;
            border: 2px solid var(--purple); color: var(--purple);
        }}
        .badge-large {{
            font-size: 0.7em; padding: 7px 22px;
        }}

        /* ── Subtitles ── */
        .slide-subtitle {{ color: var(--grey); font-size: 0.95em; margin-bottom: 20px; }}
        .slide-subtitle.small {{ font-size: 0.8em; margin-bottom: 15px; }}
        .author {{ color: var(--grey); font-size: 0.8em; margin-top: 10px; }}

        /* ── Title slide ── */
        .slide-title {{ background: var(--dark) !important; position: relative; overflow: hidden; }}
        .slide-title h1, .slide-title h2 {{ color: var(--white) !important; }}
        .slide-title .slide-subtitle {{ color: rgba(255,255,255,0.65); }}
        .slide-title .author {{ color: rgba(255,255,255,0.5); }}

        .title-bg-mesh {{
            position: absolute; inset: 0; pointer-events: none;
            background:
                radial-gradient(ellipse 60% 50% at 20% 30%, rgba(139,61,255,0.35) 0%, transparent 70%),
                radial-gradient(ellipse 50% 60% at 80% 70%, rgba(0,196,204,0.25) 0%, transparent 70%);
        }}

        .title-heading {{ position: relative; z-index: 1; }}

        /* ── Powered by Canva ── */
        .powered-by {{
            display: inline-flex; align-items: center; gap: 8px;
            margin-top: 30px; padding: 10px 20px;
            border: 1px solid rgba(255,255,255,0.2);
            border-radius: 40px; background: rgba(255,255,255,0.08);
        }}
        .powered-text {{ color: rgba(255,255,255,0.6); font-size: 0.65em; font-family: {self.FONT_BODY}; }}
        .canva-wordmark {{
            font-family: {self.FONT_DISPLAY}; font-weight: 800; font-size: 0.75em;
            background: linear-gradient(135deg, var(--purple), var(--teal));
            -webkit-background-clip: text; -webkit-text-fill-color: transparent; background-clip: text;
        }}
        .powered-by.large {{ margin-top: 35px; }}

        /* ── Section slide ── */
        .slide-section {{ position: relative; overflow: hidden; }}
        .slide-section h2 {{ color: white !important; font-size: 2em; font-weight: 800; }}
        .section-title {{ position: relative; z-index: 1; }}
        .section-sub {{ color: rgba(255,255,255,0.7) !important; position: relative; z-index: 1; font-size: 0.9em; }}
        .section-noise {{
            position: absolute; inset: 0; pointer-events: none;
            background-image: url("data:image/svg+xml,%3Csvg viewBox='0 0 200 200' xmlns='http://www.w3.org/2000/svg'%3E%3Cfilter id='n'%3E%3CfeTurbulence type='fractalNoise' baseFrequency='0.9' numOctaves='4' stitchTiles='stitch'/%3E%3C/filter%3E%3Crect width='100%25' height='100%25' filter='url(%23n)' opacity='0.04'/%3E%3C/svg%3E");
            opacity: 0.15;
        }}

        /* ── Stat grid ── */
        .stat-grid {{
            display: grid; grid-template-columns: repeat(3, 1fr);
            gap: 18px; margin: 20px auto; max-width: 1000px; width: 100%;
        }}
        .stat-grid.two {{ grid-template-columns: repeat(2, 1fr); max-width: 700px; }}
        .stat-grid.four {{ grid-template-columns: repeat(4, 1fr); max-width: 1100px; }}

        .stat-box {{
            background: var(--white); padding: 22px 18px; border-radius: 16px;
            text-align: center; border: 1px solid var(--light-grey);
            box-shadow: 0 4px 20px rgba(0,0,0,0.05);
            transition: transform 0.2s;
        }}
        .stat-box:hover {{ transform: translateY(-3px); }}
        .stat-value {{
            font-size: 2.4em; font-weight: 800; font-family: {self.FONT_DISPLAY};
            background: linear-gradient(135deg, var(--purple), var(--teal));
            -webkit-background-clip: text; -webkit-text-fill-color: transparent; background-clip: text;
            margin-bottom: 6px;
        }}
        .stat-label {{ font-size: 0.72em; color: var(--grey); font-weight: 500; }}
        .stat-sublabel {{ font-size: 0.65em; color: var(--grey); margin-top: 4px; opacity: 0.7; }}
        .stat-icon {{ font-size: 1.4em; margin-bottom: 8px; }}

        /* ── Hero stat ── */
        .hero-stat {{ margin: 20px auto; text-align: center; }}
        .hero-number {{
            font-size: 5em; font-weight: 800; font-family: {self.FONT_DISPLAY};
            background: linear-gradient(135deg, var(--purple), var(--teal));
            -webkit-background-clip: text; -webkit-text-fill-color: transparent; background-clip: text;
            line-height: 1;
        }}
        .hero-label {{ font-size: 1.1em; color: var(--grey); margin-top: 10px; font-weight: 500; }}
        .hero-sublabel {{ font-size: 0.8em; color: var(--grey); margin-top: 8px; opacity: 0.7; }}

        /* ── Flow ── */
        .flow-container {{
            display: flex; align-items: center; justify-content: center;
            flex-wrap: wrap; gap: 10px; margin: 20px auto; max-width: 1050px;
        }}
        .flow-step {{
            background: var(--white); border: 1px solid var(--light-grey);
            border-radius: 14px; padding: 18px 22px; text-align: center;
            min-width: 130px; box-shadow: 0 3px 12px rgba(0,0,0,0.05);
            flex: 1;
        }}
        .flow-step.accent {{
            background: linear-gradient(135deg, var(--purple), var(--teal));
            border-color: transparent;
        }}
        .flow-step.accent .flow-num {{ background: rgba(255,255,255,0.2); color: white; }}
        .flow-step.accent .flow-text {{ color: white; }}
        .flow-num {{
            width: 28px; height: 28px; border-radius: 50%; margin: 0 auto 10px;
            background: linear-gradient(135deg, var(--purple), var(--teal));
            color: white; font-size: 0.7em; font-weight: 700;
            display: flex; align-items: center; justify-content: center;
        }}
        .flow-text {{ font-size: 0.72em; font-weight: 500; color: var(--dark); }}
        .flow-arrow {{ font-size: 1.3em; color: var(--purple); font-weight: 700; flex-shrink: 0; }}

        /* ── Bar chart ── */
        .bar-chart {{ width: 100%; max-width: 900px; margin: 20px auto; }}
        .bar-row {{ display: flex; align-items: center; gap: 14px; margin-bottom: 14px; }}
        .bar-label {{ width: 160px; font-size: 0.72em; text-align: right; color: var(--grey); flex-shrink: 0; }}
        .bar-track {{ flex: 1; background: var(--light-grey); border-radius: 6px; height: 22px; overflow: hidden; }}
        .bar-fill {{
            height: 100%; border-radius: 6px;
            background: linear-gradient(90deg, var(--purple), var(--teal));
            transition: width 1s ease;
        }}
        .bar-value {{ width: 50px; font-size: 0.72em; font-weight: 600; color: var(--dark); }}

        /* ── Comparison ── */
        .comparison {{ display: grid; grid-template-columns: 1fr 1fr; gap: 20px; width: 100%; max-width: 900px; margin: 20px auto; }}
        .comparison-side {{ border-radius: 14px; padding: 24px; }}
        .comparison-side.left {{
            background: var(--off-white); border: 1px solid var(--light-grey);
        }}
        .comparison-side.right {{
            background: linear-gradient(135deg, var(--purple) 0%, var(--teal) 100%);
            color: white;
        }}
        .comparison-side.left h3 {{ color: var(--dark); text-align: left !important; }}
        .comparison-side.right h3 {{ color: white; text-align: left !important; }}
        .comparison-side.right li {{ color: white; }}
        .comparison-side.right li::before {{ background: rgba(255,255,255,0.7); }}

        /* ── Feature grid ── */
        .feature-grid {{ display: grid; gap: 16px; width: 100%; max-width: 1000px; margin: 20px auto; }}
        .feature-card {{
            background: var(--white); border-radius: 14px; padding: 20px;
            border: 1px solid var(--light-grey); text-align: left;
            border-left: 4px solid var(--purple); box-shadow: 0 3px 12px rgba(0,0,0,0.04);
        }}
        .feature-card h4 {{ color: var(--purple); text-align: left !important; font-size: 0.85em; }}
        .feature-card p {{ font-size: 0.72em; color: var(--grey); text-align: left !important; margin: 0; }}
        .feat-icon {{ font-size: 1.5em; margin-bottom: 8px; }}

        /* ── Quote ── */
        .quote-block {{
            max-width: 800px; margin: 20px auto;
            background: linear-gradient(135deg, rgba(139,61,255,0.06), rgba(0,196,204,0.06));
            border-left: 5px solid var(--purple); border-radius: 0 14px 14px 0;
            padding: 28px 32px; text-align: left;
        }}
        .quote-mark {{ font-size: 4em; line-height: 0.5; color: var(--purple); font-family: Georgia, serif; margin-bottom: 10px; }}
        .quote-block blockquote {{ font-size: 1em; color: var(--dark); font-style: italic; line-height: 1.6; margin: 0; }}
        .quote-author {{ margin-top: 18px; font-size: 0.75em; font-weight: 600; color: var(--purple); }}

        /* ── Phase cards ── */
        .phase-grid {{ display: flex; gap: 14px; margin: 20px auto; max-width: 1050px; width: 100%; flex-wrap: wrap; justify-content: center; }}
        .phase-card {{
            background: var(--white); border-radius: 14px; padding: 20px;
            border: 1px solid var(--light-grey); flex: 1; min-width: 170px;
            box-shadow: 0 3px 12px rgba(0,0,0,0.05); text-align: center;
        }}
        .phase-num {{
            width: 34px; height: 34px; border-radius: 50%; margin: 0 auto 8px;
            background: linear-gradient(135deg, var(--purple), var(--teal));
            color: white; font-size: 0.8em; font-weight: 700;
            display: flex; align-items: center; justify-content: center;
            font-family: {self.FONT_DISPLAY};
        }}
        .phase-label {{ font-size: 0.6em; text-transform: uppercase; letter-spacing: 1px; color: var(--teal); font-weight: 600; margin-bottom: 8px; }}
        .phase-card h4 {{ text-align: center !important; color: var(--dark); font-size: 0.82em; }}
        .phase-card p {{ font-size: 0.68em; color: var(--grey); margin: 0; text-align: center !important; }}

        /* ── Timeline ── */
        .timeline {{ position: relative; max-width: 850px; width: 100%; margin: 20px auto; }}
        .timeline::before {{
            content: ""; position: absolute; left: 50%; top: 0; bottom: 0;
            width: 2px; background: linear-gradient(to bottom, var(--purple), var(--teal));
            transform: translateX(-50%);
        }}
        .timeline-item {{ display: flex; width: 100%; margin-bottom: 20px; position: relative; }}
        .timeline-item.left {{ flex-direction: row-reverse; }}
        .timeline-content {{
            width: 44%; background: var(--white); border-radius: 12px;
            padding: 14px 18px; border: 1px solid var(--light-grey);
            box-shadow: 0 3px 10px rgba(0,0,0,0.05); text-align: left;
        }}
        .timeline-item.left .timeline-content {{ margin-left: auto; margin-right: 30px; }}
        .timeline-item.right .timeline-content {{ margin-left: 30px; }}
        .timeline-dot {{
            position: absolute; left: 50%; top: 14px;
            width: 12px; height: 12px; border-radius: 50%;
            background: linear-gradient(135deg, var(--purple), var(--teal));
            transform: translateX(-50%);
            box-shadow: 0 0 0 3px rgba(139,61,255,0.15);
        }}
        .timeline-date {{ font-size: 0.62em; font-weight: 600; color: var(--teal); margin-bottom: 4px; text-transform: uppercase; letter-spacing: 0.5px; }}
        .timeline-content h4 {{ text-align: left !important; font-size: 0.8em; color: var(--dark); }}
        .timeline-content p {{ font-size: 0.68em; color: var(--grey); margin: 0; text-align: left !important; }}

        /* ── Content list ── */
        .content-list {{ max-width: 750px; width: 100%; margin: 10px auto; }}

        /* ── Contact box ── */
        .contact-box {{
            background: rgba(255,255,255,0.08); border: 1px solid rgba(255,255,255,0.2);
            border-radius: 14px; padding: 22px 28px; margin: 20px auto;
            max-width: 520px;
        }}
        .contact-box p {{ color: rgba(255,255,255,0.85); font-size: 0.78em; margin: 8px 0; text-align: center; }}
        .contact-box strong {{ color: var(--teal); }}

        /* ── Slide numbers ── */
        .reveal .slide-number {{ font-size: 0.58em; color: var(--purple); background: transparent; opacity: 0.8; }}

        /* ── Responsive ── */
        @media (max-width: 768px) {{
            .stat-grid, .stat-grid.two, .stat-grid.four {{ grid-template-columns: 1fr; }}
            .two-col, .comparison {{ grid-template-columns: 1fr; }}
            .flow-container {{ flex-direction: column; }}
            .timeline::before {{ left: 20px; }}
            .timeline-item, .timeline-item.left {{ flex-direction: column; }}
            .timeline-content, .timeline-item.left .timeline-content {{ width: 90%; margin-left: 40px; margin-right: 0; }}
            .timeline-dot {{ left: 20px; }}
        }}
        </style>
"""

    def _build_js(self) -> str:
        return """
    <script src="https://cdn.jsdelivr.net/npm/reveal.js@4.5.0/dist/reveal.js"></script>
    <script>
        Reveal.initialize({
            controls: true,
            progress: true,
            center: true,
            hash: true,
            slideNumber: true,
            transition: 'slide',
            width: 1200,
            height: 700,
            margin: 0.1,
            minScale: 0.2,
            maxScale: 2.0,
        });
    </script>
"""

    def generate_html(self) -> str:
        slides = "\n".join(self._slides_html)
        return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{self._e(self.title)}</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.5.0/dist/reveal.css">
{self._build_css()}
</head>
<body>
<div class="reveal">
  <div class="slides">
{slides}
  </div>
</div>
{self._build_js()}
</body>
</html>"""

    def save(self, filename: str = "presentation.html", folder: str = "presentations") -> Dict:
        os.makedirs(folder, exist_ok=True)
        path = os.path.join(folder, filename)
        content = self.generate_html()
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return {
            "path": path,
            "slides": len(self._slides_html),
            "size_kb": round(len(content) / 1024, 1),
        }

    def save_pptx(self, filename: str = "presentation.pptx", folder: str = "presentations") -> Dict:
        """Export to PowerPoint. Requires: pip install python-pptx"""
        try:
            from pptx import Presentation as PptxPres
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("pip install python-pptx")

        os.makedirs(folder, exist_ok=True)
        prs = PptxPres()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        purple_rgb = RGBColor(*self.BRAND_COLORS_RGB["purple"])
        teal_rgb = RGBColor(*self.BRAND_COLORS_RGB["teal"])
        dark_rgb = RGBColor(*self.BRAND_COLORS_RGB["dark"])
        white_rgb = RGBColor(255, 255, 255)

        blank_layout = prs.slide_layouts[6]

        for slide_data in self._slides_data:
            slide = prs.slides.add_slide(blank_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.get("title", slide_data.get("quote", ""))
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = dark_rgb if slide_data.get("type") not in ("title", "section", "thank_you") else white_rgb

        path = os.path.join(folder, filename)
        prs.save(path)
        return {"path": path, "slides": len(self._slides_data)}
